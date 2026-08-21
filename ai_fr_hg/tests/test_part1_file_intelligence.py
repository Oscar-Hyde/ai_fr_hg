# Copyright (c) 2026, Ai Fr Hg and contributors
# For license information, please see license.txt

"""Part 1 File Intelligence contracts — Frappe-free unit coverage.

Covers the §6.2 format requirements, the §8 extraction result contract, the
§11 semantic grounding rules, and the archive safety policy. Everything here
runs without a bench so the guarantees can be verified in isolation; the
bench-dependent persistence paths are covered by the integration suite.
"""

from __future__ import annotations

import io
import sys
import tarfile
import types
import zipfile
from unittest import TestCase

from ai_fr_hg.ai.extraction import build_versions, detect_format, extract_bytes
from ai_fr_hg.ai.readers import get_reader, supported_extensions
from ai_fr_hg.ai.readers.archive import (
	MAX_ARCHIVE_DEPTH,
	ArchiveBudget,
	ArchiveCorruptError,
	_safe_member_path,
	is_archive,
)
from ai_fr_hg.ai.readers.code import CodeReader, language_for
from ai_fr_hg.ai.readers.container import ArchiveReader
from ai_fr_hg.ai.readers.plain import EmailReader


def _zip(members: dict[str, bytes | str]) -> bytes:
	buffer = io.BytesIO()
	with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
		for name, data in members.items():
			archive.writestr(name, data)
	return buffer.getvalue()


# ---------------------------------------------------------------------------
# §8 — the six-element extraction result contract
# ---------------------------------------------------------------------------


class TestExtractionResultContract(TestCase):
	"""§8: every extraction result carries all six mandated elements."""

	def test_evidence_contains_all_six_required_elements(self):
		outcome = extract_bytes(b"Quarterly results were strong.", "report.txt")
		evidence = outcome.evidence.as_dict()

		# 1. source identity
		self.assertTrue(evidence["provenance"]["checksum_sha256"])
		# 2. processing timestamp
		self.assertTrue(evidence["extracted_on"])
		# 3. extractor identity
		self.assertEqual(evidence["reader"], "Plain Text")
		# 4. version information
		self.assertIn("app", evidence["versions"])
		self.assertIn("reader", evidence["versions"])
		# 5. extracted content
		self.assertIn("Quarterly", outcome.result.text)
		# 6. evidence references
		self.assertIn("structure", evidence)
		self.assertIn("embedded_objects", evidence)

	def test_timestamp_is_iso_utc(self):
		from datetime import datetime

		evidence = extract_bytes(b"hello", "a.txt").evidence.as_dict()
		parsed = datetime.fromisoformat(evidence["extracted_on"])
		self.assertIsNotNone(parsed.tzinfo)

	def test_versions_identify_the_parsing_library(self):
		"""A reader backed by a library must report which version parsed."""
		from ai_fr_hg.ai.readers.office import PDFReader

		versions = build_versions(PDFReader)
		self.assertEqual(versions["library"], "pypdf")
		self.assertIn("library_version", versions)

	def test_version_reporting_never_raises(self):
		class Broken(CodeReader):
			requires = "definitely-not-installed-package-xyz"

		self.assertIsNone(Broken.library_version())
		self.assertIn("app", build_versions(Broken))

	def test_checksum_is_content_addressed(self):
		first = extract_bytes(b"same bytes", "a.txt").evidence.as_dict()
		second = extract_bytes(b"same bytes", "b.txt").evidence.as_dict()
		self.assertEqual(first["provenance"]["checksum_sha256"], second["provenance"]["checksum_sha256"])


# ---------------------------------------------------------------------------
# §6.2 — spreadsheets: formulas must not be silently discarded
# ---------------------------------------------------------------------------


class TestSpreadsheetFormulas(TestCase):
	def setUp(self):
		try:
			import openpyxl
		except ImportError:
			self.skipTest("openpyxl is not installed")

	def _workbook(self) -> bytes:
		import openpyxl

		workbook = openpyxl.Workbook()
		sheet = workbook.active
		sheet.title = "Budget"
		sheet["A1"] = "Widgets"
		sheet["B1"] = 100
		sheet["A2"] = "Total"
		sheet["B2"] = "=SUM(B1:B1)"
		buffer = io.BytesIO()
		workbook.save(buffer)
		return buffer.getvalue()

	def test_formulas_are_preserved_in_structure(self):
		result = get_reader("b.xlsx").read(self._workbook(), "b.xlsx")
		formulas = [block for block in result.structure if block["kind"] == "formula"]
		self.assertEqual(len(formulas), 1)
		self.assertEqual(formulas[0]["formula"], "=SUM(B1:B1)")
		self.assertEqual(formulas[0]["cell"], "B2")
		self.assertEqual(formulas[0]["sheet"], "Budget")

	def test_uncached_formula_values_are_reported_not_hidden(self):
		"""A workbook never opened by Excel yields blank values -- say so."""
		result = get_reader("b.xlsx").read(self._workbook(), "b.xlsx")
		self.assertEqual(result.metadata["uncached_formula_values"], 1)
		self.assertTrue(any("no cached value" in str(w) for w in result.warnings))

	def test_formula_count_is_reported(self):
		result = get_reader("b.xlsx").read(self._workbook(), "b.xlsx")
		self.assertEqual(result.metadata["formula_count"], 1)


# ---------------------------------------------------------------------------
# §6.2 — presentations: embedded content
# ---------------------------------------------------------------------------


class TestPresentationEmbeddedContent(TestCase):
	def test_pictures_are_recorded_as_embedded_objects(self):
		try:
			from pptx import Presentation
			from pptx.util import Inches
		except ImportError:
			self.skipTest("python-pptx is not installed")

		import binascii
		import struct
		import zlib

		def _png() -> bytes:
			def chunk(tag: bytes, data: bytes) -> bytes:
				body = tag + data
				return (
					struct.pack(">I", len(data)) + body + struct.pack(">I", binascii.crc32(body) & 0xFFFFFFFF)
				)

			header = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
			return (
				b"\x89PNG\r\n\x1a\n"
				+ chunk(b"IHDR", header)
				+ chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00"))
				+ chunk(b"IEND", b"")
			)

		presentation = Presentation()
		slide = presentation.slides.add_slide(presentation.slide_layouts[5])
		slide.shapes.title.text = "Review"
		slide.shapes.add_picture(io.BytesIO(_png()), Inches(1), Inches(1), Inches(1), Inches(1))
		buffer = io.BytesIO()
		presentation.save(buffer)

		result = get_reader("d.pptx").read(buffer.getvalue(), "d.pptx")
		images = [item for item in result.embedded_objects if item["kind"] == "image"]
		self.assertEqual(len(images), 1)
		self.assertEqual(images[0]["location"], "slide 1")
		self.assertEqual(result.metadata["embedded_count"], 1)

	def test_slide_order_is_preserved(self):
		try:
			from pptx import Presentation
		except ImportError:
			self.skipTest("python-pptx is not installed")

		presentation = Presentation()
		for title in ("First", "Second", "Third"):
			slide = presentation.slides.add_slide(presentation.slide_layouts[5])
			slide.shapes.title.text = title
		buffer = io.BytesIO()
		presentation.save(buffer)

		result = get_reader("d.pptx").read(buffer.getvalue(), "d.pptx")
		self.assertLess(result.text.index("First"), result.text.index("Second"))
		self.assertLess(result.text.index("Second"), result.text.index("Third"))


# ---------------------------------------------------------------------------
# §6.2 — email: conversation relationships
# ---------------------------------------------------------------------------


class TestEmailThreading(TestCase):
	MESSAGE = (
		b"From: alice@example.com\n"
		b"To: bob@example.com\n"
		b"Subject: Re: Budget\n"
		b"Message-ID: <c3@example.com>\n"
		b"In-Reply-To: <b2@example.com>\n"
		b"References: <a1@example.com> <b2@example.com>\n"
		b"Content-Type: text/plain\n\n"
		b"Agreed.\n"
	)

	def test_threading_headers_are_captured(self):
		result = EmailReader().read(self.MESSAGE, "m.eml")
		self.assertEqual(result.metadata["In-Reply-To"], "<b2@example.com>")
		self.assertTrue(result.metadata["References"])

	def test_conversation_relationship_is_reconstructable(self):
		thread = EmailReader().read(self.MESSAGE, "m.eml").metadata["thread"]
		self.assertEqual(thread["message_id"], "<c3@example.com>")
		self.assertEqual(thread["in_reply_to"], "<b2@example.com>")
		self.assertEqual(thread["references"], ["<a1@example.com>", "<b2@example.com>"])
		self.assertEqual(thread["root_message_id"], "<a1@example.com>")
		self.assertTrue(thread["is_reply"])

	def test_root_message_has_no_reply_chain(self):
		root = (
			b"From: alice@example.com\nTo: bob@example.com\nSubject: Budget\n"
			b"Message-ID: <a1@example.com>\nContent-Type: text/plain\n\nStarting point.\n"
		)
		thread = EmailReader().read(root, "root.eml").metadata["thread"]
		self.assertFalse(thread["is_reply"])
		self.assertIsNone(thread["in_reply_to"])
		self.assertEqual(thread["root_message_id"], "<a1@example.com>")


# ---------------------------------------------------------------------------
# §6.2 — source code
# ---------------------------------------------------------------------------


class TestSourceCode(TestCase):
	def test_common_languages_are_recognized(self):
		extensions = set(supported_extensions())
		for expected in ("py", "js", "ts", "java", "go", "rb", "rs", "c", "cpp", "cs", "php"):
			self.assertIn(expected, extensions, f".{expected} should be recognized")

	def test_python_structure_is_parsed_exactly(self):
		source = b'"""Doc."""\nimport os\n\n\nclass Engine:\n    def start(self):\n        pass\n'
		result = get_reader("svc.py").read(source, "svc.py")
		self.assertEqual(result.metadata["structure_fidelity"], "parsed")
		names = {block.get("name") for block in result.structure}
		self.assertIn("Engine", names)
		self.assertIn("Engine.start", names)
		self.assertIn("os", result.metadata["imports"])

	def test_indentation_is_preserved(self):
		"""Collapsing whitespace would change the meaning of Python."""
		source = b"def f():\n    if True:\n        return 1\n"
		result = get_reader("i.py").read(source, "i.py")
		self.assertIn("        return 1", result.text)

	def test_unparseable_python_degrades_honestly(self):
		result = get_reader("b.py").read(b"def broken(:\n    pass\n", "b.py")
		self.assertEqual(result.metadata["structure_fidelity"], "heuristic")
		self.assertTrue(result.warnings)

	def test_non_python_declares_heuristic_fidelity(self):
		"""A heuristic must never be presented as an authoritative parse."""
		result = get_reader("a.ts").read(b"export class Widget {}\n", "a.ts")
		self.assertEqual(result.metadata["structure_fidelity"], "heuristic")
		self.assertEqual(result.metadata["language"], "TypeScript")

	def test_module_imports_are_extracted_per_language(self):
		cases = {
			"a.ts": (b'import { thing } from "./mod";\n', "./mod"),
			"m.go": (b'import "fmt"\n', "fmt"),
			"x.java": (b"import java.util.List;\n", "java.util.List"),
			"y.c": (b"#include <stdio.h>\n", "stdio.h"),
		}
		for filename, (source, expected) in cases.items():
			result = get_reader(filename).read(source, filename)
			self.assertIn(expected, result.metadata["imports"], filename)

	def test_language_lookup(self):
		self.assertEqual(language_for("a.py"), "Python")
		self.assertIsNone(language_for("a.txt"))

	def test_detection_maps_code_to_the_code_family(self):
		self.assertEqual(detect_format(b"def f(): pass", "a.py").family, "code")


# ---------------------------------------------------------------------------
# §6.2 — archives
# ---------------------------------------------------------------------------


class TestArchiveSafety(TestCase):
	"""Protection against unsafe input is the non-negotiable archive rule."""

	def test_path_traversal_members_are_refused(self):
		result = ArchiveReader().read(_zip({"../../etc/passwd": "root", "ok.txt": "fine"}), "t.zip")
		self.assertIn("fine", result.text)
		self.assertNotIn("root", result.text)
		self.assertTrue(any("unsafe" in str(w) for w in result.warnings))

	def test_absolute_paths_are_refused(self):
		result = ArchiveReader().read(_zip({"/etc/shadow": "SHADOW_SECRET", "ok.txt": "fine"}), "t.zip")
		self.assertIn("fine", result.text)
		self.assertNotIn("SHADOW_SECRET", result.text)
		paths = {block["path"] for block in result.structure}
		self.assertNotIn("etc/shadow", paths)
		self.assertTrue(any("unsafe" in str(w) for w in result.warnings))

	def test_safe_member_path_rules(self):
		self.assertIsNone(_safe_member_path("../escape"))
		self.assertIsNone(_safe_member_path("/absolute"))
		self.assertIsNone(_safe_member_path("C:\\windows"))
		self.assertEqual(_safe_member_path("dir/file.txt"), "dir/file.txt")

	def test_compression_bomb_is_rejected(self):
		from ai_fr_hg.ai.readers.archive import ArchiveLimitError

		with self.assertRaises(ArchiveLimitError):
			ArchiveReader().read(_zip({"bomb.txt": "A" * 20_000_000}), "b.zip")

	def test_corrupt_archive_raises_a_typed_error(self):
		with self.assertRaises(ArchiveCorruptError):
			ArchiveReader().read(b"PK\x03\x04not-really-a-zip", "c.zip")

	def test_tar_links_are_never_followed(self):
		buffer = io.BytesIO()
		with tarfile.open(fileobj=buffer, mode="w") as archive:
			link = tarfile.TarInfo("evil")
			link.type = tarfile.SYMTYPE
			link.linkname = "/etc/passwd"
			archive.addfile(link)
			payload = b"safe content"
			member = tarfile.TarInfo("good.txt")
			member.size = len(payload)
			archive.addfile(member, io.BytesIO(payload))

		result = ArchiveReader().read(buffer.getvalue(), "a.tar")
		self.assertIn("safe content", result.text)
		self.assertTrue(any("link" in str(w) for w in result.warnings))


class TestArchiveProcessing(TestCase):
	def test_members_are_dispatched_to_their_own_readers(self):
		archive = _zip(
			{
				"readme.md": "# Title",
				"data.json": '{"k": "v"}',
				"code.py": "def f():\n    pass\n",
			}
		)
		result = ArchiveReader().read(archive, "bundle.zip")
		readers = {block.get("reader") for block in result.structure if block.get("read")}
		self.assertEqual(readers, {"Markdown", "JSON", "Source Code"})

	def test_nested_archives_are_processed_recursively(self):
		inner = _zip({"deep/notes.txt": "inner text"})
		outer = _zip({"nested.zip": inner, "top.txt": "outer text"})
		result = ArchiveReader().read(outer, "bundle.zip")
		self.assertIn("inner text", result.text)
		self.assertIn("outer text", result.text)

	def test_relationships_between_members_are_tracked(self):
		inner = _zip({"deep/notes.txt": "inner"})
		outer = _zip({"nested.zip": inner})
		result = ArchiveReader().read(outer, "bundle.zip")
		child = next(b for b in result.structure if b["path"] == "deep/notes.txt")
		self.assertEqual(child["parent"], "bundle.zip!nested.zip")
		self.assertEqual(child["depth"], 1)

	def test_recursion_depth_is_bounded(self):
		data = _zip({"bottom.txt": "deepest"})
		for level in range(MAX_ARCHIVE_DEPTH + 3):
			data = _zip({f"l{level}.zip": data})
		result = ArchiveReader().read(data, "deep.zip")
		self.assertLessEqual(result.metadata["max_depth_reached"], MAX_ARCHIVE_DEPTH)
		self.assertTrue(any("depth limit" in str(w) for w in result.warnings))

	def test_cumulative_member_budget_stops_the_walk(self):
		budget = ArchiveBudget(max_members=5)
		self.assertTrue(all(budget.take_member() for _ in range(5)))
		self.assertFalse(budget.take_member())
		self.assertTrue(budget.truncated)

	def test_unsupported_members_are_recorded_not_dropped(self):
		result = ArchiveReader().read(_zip({"blob.bin": b"\x00\x01", "ok.txt": "text"}), "t.zip")
		skipped = [b for b in result.structure if b.get("skipped_reason") == "unsupported_format"]
		self.assertEqual(len(skipped), 1)
		self.assertEqual(skipped[0]["path"], "blob.bin")

	def test_one_bad_member_does_not_fail_the_archive(self):
		archive = _zip({"good.txt": "readable", "broken.xlsx": b"not really xlsx"})
		result = ArchiveReader().read(archive, "t.zip")
		self.assertIn("readable", result.text)
		failed = [b for b in result.structure if b.get("skipped_reason") == "read_failed"]
		self.assertEqual(len(failed), 1)

	def test_archive_is_detected_by_extension(self):
		self.assertTrue(is_archive("a.zip"))
		self.assertTrue(is_archive("a.tar.gz"))
		self.assertFalse(is_archive("a.docx"))

	def test_office_documents_are_not_treated_as_user_archives(self):
		"""A .docx is one document, not a container of independent files."""
		self.assertFalse(is_archive("report.docx"))
		self.assertNotIsInstance(get_reader("report.docx"), ArchiveReader)


# ---------------------------------------------------------------------------
# §11 — semantic entities and relationships
# ---------------------------------------------------------------------------


def _install_frappe_stub() -> None:
	"""Minimal frappe stub so `ai.semantic` imports without a bench."""
	if "frappe" in sys.modules:
		return
	stub = types.ModuleType("frappe")

	def _throw(*args, **kwargs):
		raise Exception(args)

	stub.throw = _throw
	stub._ = lambda value: value
	utils = types.ModuleType("frappe.utils")
	utils.cint = int

	def _flt(value):
		try:
			return float(value)
		except (TypeError, ValueError):
			return 0.0

	utils.flt = _flt
	stub.utils = utils
	sys.modules["frappe"] = stub
	sys.modules["frappe.utils"] = utils


class TestSemanticGrounding(TestCase):
	"""§11: entities must be grounded, confident, and evidenced."""

	TEXT = (
		"Dr. Alice Novak works for Cyberdyne Systems in Sofia. Cyberdyne Systems is part of the Skynet Group."
	)

	@classmethod
	def setUpClass(cls):
		_install_frappe_stub()

	def _parse(self, payload, floor=50.0):
		from ai_fr_hg.ai.semantic import parse_semantic_payload

		return parse_semantic_payload(payload, self.TEXT, floor)

	def test_grounded_entities_are_accepted_with_true_offsets(self):
		out = self._parse({"entities": [{"type": "person", "value": "Alice Novak", "confidence": 95}]})
		self.assertEqual(len(out["entities"]), 1)
		entity = out["entities"][0]
		self.assertEqual(entity["entity_type"], "person")
		self.assertEqual(self.TEXT[entity["first_offset"] :][: len("Alice Novak")], "Alice Novak")

	def test_hallucinated_entities_are_rejected(self):
		"""A value the model invented is not in the document and must be dropped."""
		out = self._parse({"entities": [{"type": "person", "value": "Bob Fabricated", "confidence": 99}]})
		self.assertEqual(out["entities"], [])
		self.assertEqual(out["rejected"]["ungrounded"], 1)

	def test_low_confidence_entities_are_rejected(self):
		out = self._parse({"entities": [{"type": "person", "value": "Alice Novak", "confidence": 10}]})
		self.assertEqual(out["entities"], [])
		self.assertEqual(out["rejected"]["low_confidence"], 1)

	def test_unknown_entity_types_are_rejected(self):
		out = self._parse({"entities": [{"type": "alien", "value": "Sofia", "confidence": 90}]})
		self.assertEqual(out["entities"], [])
		self.assertEqual(out["rejected"]["invalid"], 1)

	def test_all_four_semantic_kinds_are_supported(self):
		from ai_fr_hg.ai.semantic import SEMANTIC_ENTITY_TYPES

		for kind in ("person", "organization", "location", "concept"):
			self.assertIn(kind, SEMANTIC_ENTITY_TYPES)

	def test_relationship_predicates_are_normalized(self):
		out = self._parse(
			{
				"relationships": [
					{
						"subject": "Alice Novak",
						"predicate": "works_at",
						"object": "Cyberdyne Systems",
						"evidence": "Dr. Alice Novak works for Cyberdyne Systems in Sofia.",
						"confidence": 90,
					}
				]
			}
		)
		self.assertEqual(out["relationships"][0]["relationship_type"], "works_for")

	def test_relationship_requires_grounded_evidence(self):
		out = self._parse(
			{
				"relationships": [
					{
						"subject": "Alice Novak",
						"predicate": "owns",
						"object": "Atlantis",
						"evidence": "Alice Novak owns Atlantis.",
						"confidence": 99,
					}
				]
			}
		)
		self.assertEqual(out["relationships"], [])
		self.assertEqual(out["rejected"]["ungrounded"], 1)

	def test_self_referential_relationships_are_rejected(self):
		out = self._parse(
			{
				"relationships": [
					{
						"subject": "Sofia",
						"predicate": "located_in",
						"object": "Sofia",
						"evidence": "Dr. Alice Novak works for Cyberdyne Systems in Sofia.",
						"confidence": 95,
					}
				]
			}
		)
		self.assertEqual(out["relationships"], [])

	def test_duplicate_entities_collapse(self):
		out = self._parse(
			{
				"entities": [
					{"type": "person", "value": "Alice Novak", "confidence": 95},
					{"type": "person", "value": "alice novak", "confidence": 90},
				]
			}
		)
		self.assertEqual(len(out["entities"]), 1)

	def test_malformed_payload_is_survivable(self):
		self.assertEqual(self._parse(None)["entities"], [])
		self.assertEqual(self._parse({"entities": ["not-a-dict"]})["entities"], [])

	def test_unknown_predicate_degrades_to_related_to(self):
		from ai_fr_hg.ai.semantic import normalize_relationship_type

		self.assertEqual(normalize_relationship_type("invented_predicate"), "related_to")
		self.assertEqual(normalize_relationship_type("headquartered_in"), "located_in")


class TestSemanticTypeSeparation(TestCase):
	"""A deterministic row may never claim a semantic type."""

	@classmethod
	def setUpClass(cls):
		_install_frappe_stub()

	def test_pattern_rows_cannot_carry_semantic_types(self):
		from ai_fr_hg.ai.patterns import persistable_pattern_type

		self.assertEqual(persistable_pattern_type("person", method="pattern"), "custom")
		self.assertEqual(persistable_pattern_type("person", method="semantic"), "person")

	def test_pattern_types_remain_valid_for_both(self):
		from ai_fr_hg.ai.patterns import persistable_pattern_type

		self.assertEqual(persistable_pattern_type("email", method="pattern"), "email")
		self.assertEqual(persistable_pattern_type("email", method="semantic"), "email")
